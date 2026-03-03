import sys
import subprocess

try:
    import pptx
except ImportError:
    print("python-pptx not found. Installing now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

BG = rgb('030608')
PRIMARY = rgb('1FD98A')
GOLD = rgb('F0B90B')
WHITE = rgb('FFFFFF')
SILVER = rgb('A8BFD0')
GHOST = rgb('2E4558')

def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_accent_bar(slide, color=PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_line(slide, x, y, cx, cy, color=GHOST):
    shape = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x+cx), Inches(y+cy))
    shape.line.color.rgb = color
    shape.line.width = Pt(1)

def add_text(slide, text, x, y, cx, cy, color, size, bold=False, italic=False, font_name='Georgia', align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(cx), Inches(cy))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.color.rgb = color
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    return p

def add_footer(slide, number_str):
    add_line(slide, 0.5, 7.1, 12.333, 0, GHOST)
    add_text(slide, "STABLEX  ·  HACKATHON PRESENTATION  ·  2026", 0.5, 7.15, 6, 0.3, GHOST, 9, font_name='Arial', bold=True)
    add_text(slide, number_str, 12, 7.15, 0.8, 0.3, GHOST, 9, font_name='Arial', align=PP_ALIGN.RIGHT)

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # SLIDE 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide1)
    add_accent_bar(slide1, PRIMARY)
    add_line(slide1, 6.666, 0.5, 0, 6.0, GHOST) # Vertical divider

    # Left Side
    add_text(slide1, "$92B", 0.5, 1.2, 6, 2.5, WHITE, 130, bold=True)
    add_text(slide1, "in Nigerian crypto transactions.\nJuly 2024 — June 2025.", 0.6, 3.8, 5, 1, SILVER, 19, italic=True)
    add_text(slide1, "Not a single platform handles it well.", 0.6, 4.8, 5, 0.5, PRIMARY, 15, bold=True)
    add_text(slide1, "Chainalysis  ·  Breet State of Crypto Nigeria 2025", 0.6, 6.5, 5, 0.3, GHOST, 10, italic=True)

    # Right Side
    def add_problem(slide, num, title, desc, note, y_start):
        add_text(slide, num, 7.2, y_start, 5, 0.3, PRIMARY, 10, font_name='Arial', bold=True)
        p1 = add_text(slide, title + " ", 7.2, y_start + 0.3, 5.5, 1, WHITE, 16, bold=True)
        run = p1.add_run()
        run.text = desc
        run.font.color.rgb = SILVER
        run.font.bold = False
        run.font.size = Pt(14)
        run.font.name = 'Georgia'
        add_text(slide, note, 7.2, y_start + 1.2, 5.5, 0.5, GHOST, 11, italic=True)
        if y_start < 4:
            add_line(slide, 7.2, y_start + 1.8, 5.5, 0, GHOST)

    add_problem(slide1, "01", "Fragmented.", "The average Nigerian crypto user opens 3–4 separate apps to complete one transaction. Every handoff leaks money.", "Naira lost 75% of its value since 2016.\nUsers are desperate — and trapped.", 0.5)
    add_problem(slide1, "02", "Opaque.", "Fees appear only after the transaction confirms. On every Nigerian platform. There is no way to back out.", "36% of Nigerian adults are unbanked.\nCrypto is their only financial tool.", 2.5)
    add_problem(slide1, "03", "Abandoned.", "Binance and Bybit both exited Nigeria in 2024. 22 million active users were left with nowhere to go.", "Nigeria is #2 globally in crypto adoption.\nThe demand is real. The infrastructure is not.", 4.5)

    add_footer(slide1, "01 / 03")

    # SLIDE 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2)
    add_accent_bar(slide2, GOLD)

    add_text(slide2, "WHY IT'S WORTH SOLVING", 0.5, 0.5, 5, 0.3, PRIMARY, 10, font_name='Arial', bold=True)
    add_text(slide2, "We looked at the data.\nThen we built.", 0.5, 0.8, 8, 1.5, WHITE, 52, bold=True)
    add_line(slide2, 0.5, 2.5, 12.333, 0, GHOST)

    # 3 Columns
    c1_x, c2_x, c3_x = 0.5, 4.8, 9.1
    add_line(slide2, 4.5, 2.7, 0, 3.0, GHOST)
    add_line(slide2, 8.8, 2.7, 0, 3.0, GHOST)

    def add_col(slide, x, tag, stat, desc, src, facts):
        add_text(slide, tag, x, 2.7, 3, 0.3, PRIMARY, 9, font_name='Arial', bold=True)
        add_text(slide, stat, x, 3.1, 3, 1, WHITE, 62, bold=True)
        add_text(slide, desc, x, 4.2, 3.5, 0.6, SILVER, 14)
        add_text(slide, src, x, 4.8, 3.5, 0.3, PRIMARY, 10, italic=True)
        add_line(slide, x, 5.15, 3.5, 0, GHOST)
        y = 5.25
        for f in facts:
            add_text(slide, f, x, y, 3.5, 0.3, GHOST, 11)
            y += 0.25

    add_col(slide2, c1_x, "MARKET SIZE", "#2", "Global crypto adoption rank", "Chainalysis 2024", [
        "—  $92.1B on-chain value  ·  12 months",
        "—  85% retail — not whales, everyday people",
        "—  $2.4B projected market revenue  ·  2025",
        "—  22 million users  ·  10.3% of population"
    ])

    add_col(slide2, c2_x, "USER PAIN", "43%", "of retail volume is stablecoin", "Naira hedge — not speculation", [
        "—  Naira lost 75% of value since 2016",
        "—  Inflation hit 24% in 2023",
        "—  52% of crypto users are under 30",
        "—  Fees hidden until after confirmation"
    ])

    add_col(slide2, c3_x, "MARKET GAP", "0", "local platforms with full stack", "Busha · Quidax · Patricia · Yellow Card", [
        "—  No local app: BTC + ETH + SOL + TRON",
        "—  Binance & Bybit both exited 2024",
        "—  No custodial wallet + NGN on/off ramp",
        "—  This is the gap StableX was built for"
    ])

    # Gold Bar
    gold_bar = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.8))
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = GOLD
    gold_bar.line.fill.background()

    add_text(slide2, "THE WINDOW IS NOW.", 0.8, 6.4, 4, 0.5, rgb('000000'), 18, bold=True)
    add_line(slide2, 4.0, 6.45, 0, 0.3, rgb('000000'))
    add_text(slide2, "Binance & Bybit exited. Regulation arrived. 22M users need a home. We built it.", 4.25, 6.43, 8.5, 0.5, rgb('000000'), 14, bold=True)

    add_footer(slide2, "02 / 03")

    # SLIDE 3
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide3)
    add_accent_bar(slide3, PRIMARY)

    add_line(slide3, 6.666, 0.5, 0, 6.0, GHOST)

    # Left Side
    add_text(slide3, "WHAT WE BUILT", 0.5, 0.5, 5, 0.3, PRIMARY, 10, font_name='Arial', bold=True)
    add_text(slide3, "StableX", 0.5, 0.8, 6, 1.5, WHITE, 72, bold=True)
    add_text(slide3, "Nigeria's All-In-One\nCrypto Finance Platform", 0.6, 2.4, 5, 1, SILVER, 18, italic=True)

    add_text(slide3, "SHIPPED AS A WORKING MVP:", 0.6, 3.6, 5, 0.3, PRIMARY, 10, font_name='Arial', bold=True)

    y_feat = 4.0
    feats_left = [
        "—  4 live blockchains: BTC, ETH, SOL, TRON",
        "—  Real sweep bots: automated fund management",
        "—  KoraPay + Interswitch NGN on/off ramp",
        "—  Custodial wallets with AES-256 encryption"
    ]
    feats_right = [
        "—  Staking & Earn: 8% APY yield product",
        "—  Merchant & Developer API with webhooks",
        "—  Gift cards via Reloadly integration",
        "—  JWT + OTP + Google OAuth full auth stack"
    ]

    for i, f in enumerate(feats_left):
        add_text(slide3, f, 0.6, y_feat + i*0.4, 2.8, 0.4, SILVER, 10.5)

    for i, f in enumerate(feats_right):
        add_text(slide3, f, 3.6, y_feat + i*0.4, 2.8, 0.4, SILVER, 10.5)

    # Live Now box
    box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.8), Inches(5.5), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = BG
    box.line.color.rgb = PRIMARY
    box.line.width = Pt(1.5)
    add_text(slide3, "LIVE NOW →  stablexv1-production.up.railway.app", 0.6, 6.0, 5.5, 0.4, PRIMARY, 12, bold=True, align=PP_ALIGN.CENTER)

    # Right Side
    add_text(slide3, "PRODUCT DEMO", 7.2, 0.5, 5, 0.3, PRIMARY, 10, font_name='Arial', bold=True)
    
    vid_box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.2), Inches(0.9), Inches(5.5), Inches(5.5))
    vid_box.fill.solid()
    vid_box.fill.fore_color.rgb = rgb('0A1014')
    vid_box.line.color.rgb = GHOST
    vid_box.line.width = Pt(1)

    # Thin emerald line across top of video box
    add_line(slide3, 7.2, 0.9, 5.5, 0, PRIMARY)
    slide3.shapes[-1].line.width = Pt(2)

    # Play circle
    circle = slide3.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.55), Inches(3.25), Inches(0.8), Inches(0.8))
    circle.fill.background()
    circle.line.color.rgb = PRIMARY
    circle.line.width = Pt(1.5)

    # Play arrow (triangle)
    arrow = slide3.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(9.85), Inches(3.45), Inches(0.25), Inches(0.4))
    arrow.rotation = 90
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = PRIMARY
    arrow.line.fill.background()

    add_text(slide3, "INSERT DEMO VIDEO  ·  4 MINUTES MAX", 7.2, 4.3, 5.5, 0.3, GHOST, 10, align=PP_ALIGN.CENTER, font_name='Arial')

    add_footer(slide3, "03 / 03")

    prs.save('StableX_Hackathon.pptx')
    print("Successfully generated StableX_Hackathon.pptx")

if __name__ == "__main__":
    create_deck()
