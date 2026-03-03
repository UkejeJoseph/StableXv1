import sys
import subprocess

try:
    import pptx
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

# Modern Dark Theme
BG_COLOR = rgb('0F172A')       # Slate 900
CARD_BG = rgb('1E293B')        # Slate 800
TEXT_MAIN = rgb('F8FAFC')      # Slate 50
TEXT_MUTED = rgb('94A3B8')     # Slate 400
BRAND_COLOR = rgb('10B981')    # Emerald 500
BRAND_GLOW = rgb('064E3B')     # Emerald 900
GOLD = rgb('F59E0B')           # Amber 500

def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title, subtitle):
    # Emerald accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.75), Inches(0.12), Inches(0.85))
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND_COLOR
    line.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(0.6), Inches(10), Inches(0.6))
    p = txBox.text_frame.paragraphs[0]
    p.text = title
    p.font.name = 'Arial'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(10), Inches(0.5))
    p2 = txBox2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.name = 'Arial'
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_MUTED

def add_card(slide, x, y, w, h, title, text, title_color=BRAND_COLOR):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = rgb('334155') # Slate 700
    card.line.width = Pt(1.5)

    # Card Title
    tx1 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(w - 0.4), Inches(0.4))
    p1 = tx1.text_frame.paragraphs[0]
    p1.text = title
    p1.font.name = 'Arial'
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.font.color.rgb = title_color

    # Card Text
    tx2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(w - 0.4), Inches(h - 0.9))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = text
    p2.font.name = 'Arial'
    p2.font.size = Pt(15)
    p2.font.color.rgb = TEXT_MAIN
    p2.line_spacing = 1.3

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ==========================================
    # SLIDE 1: THE PROBLEM
    # ==========================================
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s1)
    add_header(s1, "The Nigerian Crypto Crisis", "Survival, Not Speculation")
    
    # 3-Card layout
    add_card(s1, 0.8, 2.2, 3.6, 4.5, "Massive Scale", 
             "Nigeria has tens of millions of active crypto users moving billions of dollars through the market annually. Yet, no single platform exists that handles everything in one place.")
    
    add_card(s1, 4.8, 2.2, 3.6, 4.5, "Extreme Friction", 
             "Users juggle three to four separate apps just to complete one simple transaction. Fees are completely hidden until it is too late to cancel, silently draining wealth.")

    add_card(s1, 8.8, 2.2, 3.6, 4.5, "Total Abandonment", 
             "The biggest global platforms abandoned Nigerian users in 2024. Meanwhile, the Naira is collapsing. For Nigerians, crypto is not speculation—it is absolute financial survival.\n\nThis problem is real, urgent, and completely unsolved.", title_color=GOLD)


    # ==========================================
    # SLIDE 2: VALIDATION
    # ==========================================
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s2)
    add_header(s2, "Qualifying the Problem", "Backed by market data and true user research")

    # Top two wide cards
    add_card(s2, 0.8, 2.2, 5.7, 2.0, "#2 Globally in Adoption", 
             "Nigeria's billions of dollars in volume are driven purely by ordinary people trying to survive inflation, not institutions.")
    
    add_card(s2, 6.9, 2.2, 5.5, 2.0, "100% User Frustration", 
             "We spoke to real users. They confirmed the daily nightmare of toggling multiple apps for a single transaction. No local competitor offers a complete solution.")

    # Giant verdict bar at the bottom
    verdict_bg = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.5), Inches(11.6), Inches(2.2))
    verdict_bg.fill.solid()
    verdict_bg.fill.fore_color.rgb = BRAND_GLOW
    verdict_bg.line.color.rgb = BRAND_COLOR
    verdict_bg.line.width = Pt(2)

    v_tx = s2.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.2), Inches(1.5))
    tf_v = v_tx.text_frame
    tf_v.word_wrap = True
    vp = tf_v.paragraphs[0]
    vp.text = "THE WINDOW OF OPPORTUNITY IS OPEN RIGHT NOW."
    vp.font.name = 'Arial'
    vp.font.size = Pt(24)
    vp.font.bold = True
    vp.font.color.rgb = BRAND_COLOR
    vp.alignment = PP_ALIGN.CENTER

    vp2 = tf_v.add_paragraph()
    vp2.text = "The biggest global platforms left the market. Millions of users are actively looking for something better today."
    vp2.font.name = 'Arial'
    vp2.font.size = Pt(20)
    vp2.font.color.rgb = TEXT_MAIN
    vp2.alignment = PP_ALIGN.CENTER


    # ==========================================
    # SLIDE 3: WORKING PRODUCT / DEMO
    # ==========================================
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s3)
    add_header(s3, "StableX: The All-In-One Platform", "Not a prototype. A fully working platform live right now.")

    # Left Side: Features
    features = [
        "4 Supported Blockchains: BTC, ETH, SOL, TRON",
        "Seamless deposits & withdrawals directly in Naira",
        "Instant Swapping, Staking Yields, and Gift Cards",
        "Developer Merchant API for business integration",
        "Secure, bank-grade AES-256 encrypted wallets",
        "",
        "The fact that we shipped a complete working product is what separates us from every other team here."
    ]

    ftx = s3.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(6.0), Inches(4.5))
    tf = ftx.text_frame
    tf.word_wrap = True
    for i, feature in enumerate(features):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        if i == len(features) - 1:
            p.text = feature
            p.font.color.rgb = BRAND_COLOR
            p.font.bold = True
        else:
            if feature:
                p.text = "✦  " + feature
            else:
                p.text = feature
            p.font.color.rgb = TEXT_MAIN
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.line_spacing = 1.4

    # Right Side: Video Placeholder
    vid_ph = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(2.2), Inches(5.2), Inches(4.5))
    vid_ph.fill.solid()
    vid_ph.fill.fore_color.rgb = rgb('000000')
    vid_ph.line.color.rgb = rgb('334155')
    vid_ph.line.width = Pt(2)

    v_text = s3.shapes.add_textbox(Inches(7.2), Inches(4.0), Inches(5.2), Inches(1.0))
    vpText = v_text.text_frame.paragraphs[0]
    vpText.text = "▶ 4-MINUTE DEMO VIDEO"
    vpText.font.name = 'Arial'
    vpText.font.size = Pt(20)
    vpText.font.bold = True
    vpText.font.color.rgb = TEXT_MUTED
    vpText.alignment = PP_ALIGN.CENTER


    # Save
    prs.save('StableX_Demo_Day.pptx')
    print("Successfully generated StableX_Demo_Day.pptx")

if __name__ == "__main__":
    create_deck()
